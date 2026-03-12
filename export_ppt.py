"""PPT 생성 모듈 — CEO 보고용 프리미엄 디자인"""
import io, os, json
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm

DATA_DIR = os.path.join(os.path.dirname(__file__), "data")
KR_FONT = "맑은 고딕"

# ── 색상 팔레트 (딥블루 모던) ──
C_PRIMARY = RGBColor(0x1e, 0x3a, 0x8a)   # 딥 블루
C_ACCENT  = RGBColor(0x3b, 0x82, 0xf6)   # 액센트 블루
C_PURPLE  = RGBColor(0x63, 0x66, 0xf1)   # 퍼플
C_DARK    = RGBColor(0x0f, 0x17, 0x2a)
C_TEXT    = RGBColor(0x33, 0x41, 0x55)
C_GRAY    = RGBColor(0x64, 0x74, 0x8b)
C_LGRAY   = RGBColor(0x94, 0xa3, 0xb8)
C_WHITE   = RGBColor(0xff, 0xff, 0xff)
C_LIGHT   = RGBColor(0xf1, 0xf5, 0xf9)
C_LIGHTER = RGBColor(0xf8, 0xfa, 0xfc)
C_BORDER  = RGBColor(0xe2, 0xe8, 0xf0)
C_GREEN   = RGBColor(0x05, 0x96, 0x69)
C_ORANGE  = RGBColor(0xd9, 0x77, 0x06)
C_RED     = RGBColor(0xdc, 0x26, 0x26)
C_HEADER_SUB = RGBColor(0x93, 0xc5, 0xfd)
C_BADGE   = RGBColor(0xfb, 0xbf, 0x24)


def _fmt(n, suffix=""):
    try:
        n = int(n or 0)
        if n >= 10_000_000: return f"{n/10_000_000:.1f}천만{suffix}"
        if n >= 10_000: return f"{n/10_000:.1f}만{suffix}"
        return f"{n:,}{suffix}"
    except: return str(n or "-")

def _fmt_rate(r):
    try: return f"{float(r or 0):.2f}%"
    except: return "-"

def _er_color(er):
    try:
        v = float(er or 0)
        if v >= 5: return C_GREEN
        if v >= 2: return C_ACCENT
        return C_ORANGE
    except: return C_GRAY


def _tb(slide, text, left, top, w, h, size=10, bold=False, color=C_DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(w), Cm(h))
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(text)
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = KR_FONT
    return box

def _rect(slide, left, top, w, h, fill_color, line=False):
    shape = slide.shapes.add_shape(1, Cm(left), Cm(top), Cm(w), Cm(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line:
        shape.line.color.rgb = C_BORDER
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape

def _stat_card(slide, label, value, x, y, w, h, value_color=C_PRIMARY):
    _rect(slide, x, y, w, h, C_WHITE, line=True)
    _tb(slide, value, x, y + 0.15, w, h * 0.55, size=16, bold=True, color=value_color, align=PP_ALIGN.CENTER)
    _tb(slide, label, x, y + h * 0.55, w, h * 0.38, size=8, color=C_GRAY, align=PP_ALIGN.CENTER)


def _add_scorecard_slide(prs, inf, manual):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = 33.87, 19.05

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_LIGHTER

    username = inf.get("username", "")
    full_name = inf.get("full_name", "")

    # ═══ 1. 헤더 바 (딥블루) ═══
    _rect(slide, 0, 0, W, 3.2, C_PRIMARY)

    # 프로필 사진
    import requests as _req
    pk = inf.get("pk", "")
    _pic_urls = []
    if pk:
        sb_url = os.environ.get("SUPABASE_URL", "https://ysqnixgdpltguatvjjcb.supabase.co")
        _pic_urls.append(f"{sb_url}/storage/v1/object/public/profile-pics/{pk}.jpg")
    _local = inf.get("profile_pic_local", "") or ""
    if _local.startswith("http"): _pic_urls.append(_local)
    _cdn = inf.get("profile_pic_url", "") or ""
    if _cdn: _pic_urls.append(_cdn)
    for _purl in _pic_urls:
        try:
            _pr = _req.get(_purl, timeout=5)
            if _pr.status_code == 200 and len(_pr.content) > 500:
                slide.shapes.add_picture(io.BytesIO(_pr.content), Cm(1.2), Cm(0.35), Cm(2.5), Cm(2.5))
                break
        except: continue

    # 이름 + 배지
    _tb(slide, f"@{username}", 4.2, 0.25, 16, 1.0, size=22, bold=True, color=C_WHITE)
    sub = full_name
    cat = inf.get("category") or manual.get("main_category") or ""
    if cat: sub += f"  ·  {cat}"
    if inf.get("is_verified"): sub += "  ·  인증"
    if manual.get("can_live"): sub += "  ·  라이브"
    _tb(slide, sub, 4.2, 1.35, 16, 0.6, size=10, color=C_HEADER_SUB)
    _tb(slide, f"instagram.com/{username}", 4.2, 2.05, 16, 0.5, size=8, color=C_HEADER_SUB)

    # 오른쪽 로고
    _tb(slide, "SUPERTAG", W - 5.5, 0.3, 5, 0.8, size=13, bold=True, color=C_WHITE, align=PP_ALIGN.RIGHT)
    _tb(slide, "Influencer Scorecard", W - 5.5, 1.15, 5, 0.5, size=8, color=C_HEADER_SUB, align=PP_ALIGN.RIGHT)
    _tb(slide, datetime.now().strftime('%Y.%m.%d'), W - 5.5, 1.7, 5, 0.5, size=9, color=C_HEADER_SUB, align=PP_ALIGN.RIGHT)

    # ═══ 2. 핵심 지표 6개 ═══
    y = 3.6
    cw = (W - 1.4) / 6
    er = inf.get("engagement_rate", 0)
    stats = [
        ("팔로워", _fmt(inf.get("follower_count")), C_PRIMARY),
        ("팔로잉", _fmt(inf.get("following_count")), C_TEXT),
        ("총 게시물", _fmt(inf.get("media_count")), C_TEXT),
        ("참여율", _fmt_rate(er), _er_color(er)),
        ("릴스 비율", _fmt_rate(inf.get("reels_ratio")), C_ACCENT),
        ("협찬 비율", _fmt_rate(inf.get("sponsored_ratio")), C_ORANGE),
    ]
    for i, (lbl, val, clr) in enumerate(stats):
        _stat_card(slide, lbl, val, 0.5 + i * (cw + 0.08), y, cw, 1.9, clr)

    # ═══ 3. 피드/릴스 성과 비교 ═══
    y2 = 5.9
    sw = 15.5
    _rect(slide, 0.5, y2, sw, 0.5, C_PRIMARY)
    _tb(slide, "피드 / 릴스 성과 비교", 0.7, y2 + 0.03, 10, 0.45, size=9, bold=True, color=C_WHITE)

    pcw = sw / 4
    hy = y2 + 0.55
    for i, h in enumerate(["", "평균 좋아요", "평균 댓글", "평균 조회수"]):
        _rect(slide, 0.5 + i * pcw, hy, pcw, 0.55, C_LIGHT, line=True)
        _tb(slide, h, 0.5 + i * pcw, hy + 0.05, pcw, 0.45, size=7.5, bold=True, color=C_GRAY, align=PP_ALIGN.CENTER)

    # 피드 행 (연한 파랑 배경)
    fy = hy + 0.6
    _rect(slide, 0.5, fy, pcw, 0.6, RGBColor(0xef, 0xf6, 0xff), line=True)
    _tb(slide, "피드 (Feed)", 0.7, fy + 0.08, pcw - 0.4, 0.45, size=8, bold=True, color=C_ACCENT)
    for i, v in enumerate([_fmt(inf.get("avg_feed_likes")), _fmt(inf.get("avg_feed_comments")), "-"]):
        _rect(slide, 0.5 + pcw * (i + 1), fy, pcw, 0.6, RGBColor(0xef, 0xf6, 0xff), line=True)
        _tb(slide, v, 0.5 + pcw * (i + 1), fy + 0.08, pcw, 0.45, size=9, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

    # 릴스 행 (연한 보라 배경)
    ry = fy + 0.65
    _rect(slide, 0.5, ry, pcw, 0.6, RGBColor(0xf5, 0xf3, 0xff), line=True)
    _tb(slide, "릴스 (Reels)", 0.7, ry + 0.08, pcw - 0.4, 0.45, size=8, bold=True, color=C_PURPLE)
    for i, v in enumerate([_fmt(inf.get("avg_reel_likes")), _fmt(inf.get("avg_reel_comments")), _fmt(inf.get("avg_reel_views"))]):
        _rect(slide, 0.5 + pcw * (i + 1), ry, pcw, 0.6, RGBColor(0xf5, 0xf3, 0xff), line=True)
        clr = C_PURPLE if i == 2 else C_DARK
        _tb(slide, v, 0.5 + pcw * (i + 1), ry + 0.08, pcw, 0.45, size=9, bold=True, color=clr, align=PP_ALIGN.CENTER)

    # ═══ 4. 계정 정보 (오른쪽) ═══
    rx = 16.5
    rw = W - rx - 0.5
    _rect(slide, rx, y2, rw, 0.5, C_PRIMARY)
    _tb(slide, "계정 정보 / 단가", rx + 0.2, y2 + 0.03, 10, 0.45, size=9, bold=True, color=C_WHITE)

    info = [
        ("카테고리", inf.get("category") or manual.get("main_category") or "-"),
        ("업로드 빈도", inf.get("upload_frequency") or "-"),
        ("마지막 게시", inf.get("last_post_date") or "-"),
        ("프로필 링크", (inf.get("external_url") or "-")[:30]),
        ("이메일", manual.get("contact_email") or inf.get("public_email") or "-"),
        ("피드 단가", f"{manual.get('feed_price') or 0:,}만원" if manual.get("feed_price") else "-"),
        ("릴스 단가", f"{manual.get('reel_price') or 0:,}만원" if manual.get("reel_price") else "-"),
        ("협업 유형", manual.get("collab_types") or "-"),
    ]
    iy = y2 + 0.55
    for i, (lbl, val) in enumerate(info):
        bg = C_WHITE if i % 2 == 0 else C_LIGHTER
        _rect(slide, rx, iy + i * 0.55, rw, 0.55, bg, line=True)
        _tb(slide, lbl, rx + 0.2, iy + i * 0.55 + 0.05, 3.5, 0.45, size=7.5, bold=True, color=C_GRAY)
        _tb(slide, val, rx + 3.8, iy + i * 0.55 + 0.05, rw - 4, 0.45, size=7.5, color=C_DARK)

    # ═══ 5. 인기 게시물 ═══
    y3 = 10.0
    tw = W - 1
    _rect(slide, 0.5, y3, tw, 0.5, C_PRIMARY)
    _tb(slide, "인기 게시물 TOP", 0.7, y3 + 0.03, 10, 0.45, size=9, bold=True, color=C_WHITE)

    pcols = [2.5, 7, 4.5, 4.5, 4.5]
    pheaders = ["유형", "게시물 코드", "좋아요", "댓글", "조회수"]
    px = 0.5
    ph_y = y3 + 0.55
    for i, (h, pw) in enumerate(zip(pheaders, pcols)):
        _rect(slide, px, ph_y, pw, 0.5, C_LIGHT, line=True)
        _tb(slide, h, px, ph_y + 0.05, pw, 0.4, size=7.5, bold=True, color=C_GRAY, align=PP_ALIGN.CENTER)
        px += pw

    top_likes = inf.get("top_posts_likes", [])
    top_reels = inf.get("top_reels_views", [])
    all_posts = []
    for p in (top_likes or [])[:3]:
        if isinstance(p, dict):
            url = p.get("url", "")
            code = url.split("/p/")[-1].rstrip("/")[:15] if "/p/" in url else "-"
            all_posts.append(("피드", code, _fmt(p.get("likes", 0)), _fmt(p.get("comments", 0)), "-"))
    for p in (top_reels or [])[:3]:
        if isinstance(p, dict):
            url = p.get("url", "")
            code = url.split("/p/")[-1].rstrip("/")[:15] if "/p/" in url else "-"
            all_posts.append(("릴스", code, _fmt(p.get("likes", 0)), "-", _fmt(p.get("views", 0))))

    for ri, row in enumerate(all_posts[:6]):
        pr_y = ph_y + 0.55 + ri * 0.48
        bg = C_WHITE if ri % 2 == 0 else C_LIGHTER
        px = 0.5
        for ci, (val, pw) in enumerate(zip(row, pcols)):
            _rect(slide, px, pr_y, pw, 0.48, bg, line=True)
            if ci == 0:
                clr = C_ACCENT if val == "피드" else C_PURPLE
            elif ci == 4 and val != "-":
                clr = C_PURPLE
            else:
                clr = C_DARK
            _tb(slide, val, px, pr_y + 0.04, pw, 0.4, size=7.5, color=clr,
                align=PP_ALIGN.CENTER, bold=(ci == 0))
            px += pw

    # ═══ 6. 해시태그 ═══
    y4 = 13.6
    htags_raw = inf.get("top_hashtags", "[]")
    try: htags = json.loads(htags_raw) if isinstance(htags_raw, str) else (htags_raw or [])
    except: htags = []
    if htags:
        _rect(slide, 0.5, y4, tw, 0.45, C_PRIMARY)
        _tb(slide, "자주 사용하는 해시태그", 0.7, y4 + 0.02, 10, 0.4, size=8.5, bold=True, color=C_WHITE)
        tag_text = "  ".join([f"#{h['tag']}({h['count']})" if isinstance(h, dict) else f"#{h}" for h in htags[:15]])
        _rect(slide, 0.5, y4 + 0.5, tw, 0.7, C_WHITE, line=True)
        _tb(slide, tag_text, 0.7, y4 + 0.55, tw - 0.4, 0.6, size=8, color=C_PURPLE)

    # ═══ 7. 메모 + 바이오 ═══
    my = 15.0
    notes = manual.get("notes", "")
    brands = manual.get("past_brands", "")
    bio = inf.get("biography", "")
    has_info = notes or brands or bio
    if has_info:
        _rect(slide, 0.5, my, tw, 1.5, C_WHITE, line=True)
        ty = my + 0.08
        if bio:
            _tb(slide, f"바이오: {bio[:120]}", 0.7, ty, tw - 0.4, 0.4, size=7, color=C_TEXT)
            ty += 0.45
        if notes:
            _tb(slide, f"메모: {notes}", 0.7, ty, tw - 0.4, 0.4, size=7, color=C_GRAY)
            ty += 0.4
        if brands:
            _tb(slide, f"협업 브랜드: {brands}", 0.7, ty, tw - 0.4, 0.4, size=7, color=C_GRAY)

    # ═══ 8. 푸터 ═══
    _rect(slide, 0, H - 0.55, W, 0.55, C_PRIMARY)
    _tb(slide, f"SuperTag  |  Confidential  |  {datetime.now().strftime('%Y-%m-%d %H:%M')}",
        0.5, H - 0.5, W - 1, 0.4, size=7, color=C_HEADER_SUB, align=PP_ALIGN.RIGHT)


def export_single_ppt(inf: dict, manual: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = Cm(33.87)
    prs.slide_height = Cm(19.05)
    _add_scorecard_slide(prs, inf, manual)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def export_multi_ppt(inf_list: list) -> bytes:
    prs = Presentation()
    prs.slide_width = Cm(33.87)
    prs.slide_height = Cm(19.05)
    for inf, manual in inf_list:
        _add_scorecard_slide(prs, inf, manual)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def export_list_ppt(inf_list: list) -> bytes:
    prs = Presentation()
    prs.slide_width = Cm(33.87)
    prs.slide_height = Cm(19.05)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = 33.87, 19.05

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_LIGHTER

    _rect(slide, 0, 0, W, 2.2, C_PRIMARY)
    _tb(slide, "인플루언서 비교 리스트", 0.8, 0.3, 20, 1, size=20, bold=True, color=C_WHITE)
    _tb(slide, f"총 {len(inf_list)}명  |  {datetime.now().strftime('%Y-%m-%d %H:%M')}  |  SuperTag",
        0.8, 1.3, 20, 0.6, size=9, color=C_HEADER_SUB)

    rows = len(inf_list) + 1
    cols = 11
    left, top = Cm(0.3), Cm(2.5)
    width, height = Cm(W - 0.6), Cm(H - 3.2)

    tbl = slide.shapes.add_table(rows, cols, left, top, width, height).table
    headers = ["계정", "팔로워", "참여율", "피드 평균\n좋아요", "피드 평균\n댓글",
               "릴스 평균\n조회수", "릴스 평균\n좋아요", "게시물", "카테고리", "피드단가", "릴스단가"]

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0] if p.runs else p.add_run()
        r.font.bold = True
        r.font.size = Pt(7.5)
        r.font.color.rgb = C_WHITE
        r.font.name = KR_FONT

    for i, (inf, manual) in enumerate(inf_list):
        vals = [
            f"@{inf.get('username', '')}",
            _fmt(inf.get("follower_count")),
            _fmt_rate(inf.get("engagement_rate")),
            _fmt(inf.get("avg_feed_likes")),
            _fmt(inf.get("avg_feed_comments")),
            _fmt(inf.get("avg_reel_views")),
            _fmt(inf.get("avg_reel_likes")),
            _fmt(inf.get("media_count")),
            (inf.get("category") or manual.get("main_category") or "-")[:10],
            f"{manual.get('feed_price') or 0}만" if manual.get("feed_price") else "-",
            f"{manual.get('reel_price') or 0}만" if manual.get("reel_price") else "-",
        ]
        bg = C_LIGHTER if i % 2 == 0 else C_WHITE
        for j, val in enumerate(vals):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            r = p.runs[0] if p.runs else p.add_run()
            r.font.size = Pt(7.5)
            r.font.color.rgb = C_DARK
            r.font.name = KR_FONT

    _rect(slide, 0, H - 0.5, W, 0.5, C_PRIMARY)
    _tb(slide, "SuperTag  |  Confidential", 0.5, H - 0.45, W - 1, 0.4,
        size=7, color=C_HEADER_SUB, align=PP_ALIGN.RIGHT)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
